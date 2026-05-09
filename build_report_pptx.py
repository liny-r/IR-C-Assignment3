from pathlib import Path
from PIL import Image as PILImage

BASE    = Path("/Users/yueqilin/Desktop/MTH9877 IR/IR&C Assignment3")
IMG_DIR = BASE / "processed"
OUT_DIR = BASE / "processed"

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
def img(name):
    p = IMG_DIR / name
    return str(p) if p.exists() else None

# ─────────────────────────────────────────────────────────────────────────────
# 1. PDF REPORT  (reportlab)
# ─────────────────────────────────────────────────────────────────────────────
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor, black, white
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak,
    Table, TableStyle, HRFlowable,
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

W, H = letter
styles = getSampleStyleSheet()

# Madrid theme colors
BLUE   = HexColor("#1F3963")   # deep navy
ACCENT = HexColor("#C00000")   # crimson red
GOLD   = HexColor("#FFC000")   # gold
TEAL   = HexColor("#00567A")   # teal
RED    = HexColor("#C00000")
LGREY  = HexColor("#F2F2F2")
MGREY  = HexColor("#595959")

title_style = ParagraphStyle("Title2", parent=styles["Title"],
    fontSize=22, textColor=BLUE, spaceAfter=6, alignment=TA_CENTER)
h1_style = ParagraphStyle("H1", parent=styles["Heading1"],
    fontSize=14, textColor=white, spaceAfter=0,
    backColor=BLUE, borderPad=5, leading=18)
h2_style = ParagraphStyle("H2", parent=styles["Heading2"],
    fontSize=12, textColor=BLUE, spaceAfter=4, spaceBefore=8)
body_style = ParagraphStyle("Body2", parent=styles["Normal"],
    fontSize=9.5, leading=14, spaceAfter=4, alignment=TA_JUSTIFY)
bullet_style = ParagraphStyle("Bullet", parent=body_style,
    leftIndent=14, bulletIndent=4, spaceAfter=2)
caption_style = ParagraphStyle("Caption", parent=styles["Normal"],
    fontSize=8, textColor=HexColor("#595959"),
    alignment=TA_CENTER, spaceAfter=8)
small_style = ParagraphStyle("Small", parent=body_style, fontSize=8.5)
formula_style = ParagraphStyle("Formula", parent=styles["Normal"],
    fontSize=10, leading=16, fontName="Courier",
    backColor=HexColor("#EEF2F7"), leftIndent=20, rightIndent=20,
    borderPad=6, spaceAfter=6, spaceBefore=4)

def h1(text):
    return [
        Paragraph(f"<font color='white'>&nbsp;&nbsp;{text}</font>", h1_style),
        HRFlowable(width="100%", thickness=3, color=GOLD, spaceAfter=6),
    ]

def h2(text):
    return Paragraph(text, h2_style)

def p(text):
    return Paragraph(text, body_style)

def b(text):
    return Paragraph(f"• {text}", bullet_style)

def caption(text):
    return Paragraph(f"<i>{text}</i>", caption_style)

def sp(n=6):
    return Spacer(1, n)

def hr():
    return HRFlowable(width="100%", thickness=0.5, color=LGREY, spaceAfter=4)

def formula(text):
    return Paragraph(text, formula_style)

def add_image(path, max_w=6.5*inch, max_h=3.5*inch):
    if not path:
        return sp(4)
    try:
        with PILImage.open(path) as im:
            pw, ph = im.size
        scale = min(max_w / pw, max_h / ph)
        return Image(path, width=pw*scale, height=ph*scale)
    except Exception:
        return sp(4)

def two_col_images(p1, p2, max_w=3.1*inch, max_h=2.8*inch):
    i1 = add_image(p1, max_w, max_h)
    i2 = add_image(p2, max_w, max_h)
    t = Table([[i1, i2]], colWidths=[3.25*inch, 3.25*inch])
    t.setStyle(TableStyle([("ALIGN", (0,0), (-1,-1), "CENTER"),
                            ("VALIGN", (0,0), (-1,-1), "MIDDLE")]))
    return t

# ── Build story ───────────────────────────────────────────────────────────────
story = []

# Cover
story += [
    sp(30),
    Paragraph("MTH9877 — Interest Rates & Credit", ParagraphStyle("Cover1",
        parent=styles["Normal"], fontSize=13, textColor=ACCENT, alignment=TA_CENTER)),
    sp(6),
    Paragraph("Assignment 3 — Part E: Extensions", title_style),
    sp(8),
    HRFlowable(width="60%", thickness=2, color=ACCENT, spaceAfter=8),
    Paragraph("Mortgage Prepayment &amp; Default Survival Models", ParagraphStyle(
        "Sub", parent=styles["Normal"], fontSize=11, alignment=TA_CENTER, textColor=BLUE)),
    sp(24),
    Paragraph("Rose Lin · Baruch MFE · May 2026", ParagraphStyle(
        "Auth", parent=styles["Normal"], fontSize=10, alignment=TA_CENTER,
        textColor=HexColor("#7f8c8d"))),
    PageBreak(),
]

# ── E(i): Competing Risks ─────────────────────────────────────────────────────
story += h1("E(i) — Competing Risks Framework"); story += [sp(4)]

story += [h2("(a) Exploratory Data Analysis  —  Full 34M-Loan Dataset"), sp(2)]
story += [
    p("The Freddie Mac dataset contains 34 million 30-year fixed-rate mortgages originated "
      "1999–2023. Loans exit by prepayment (~57%), remain active / censored (~41%), or "
      "default (~2%)."),
    sp(4),
    add_image(img("E1_eda_duration.png"), max_h=2.8*inch),
    caption("Fig 1. Duration distributions by event type. "
            "Prepaid loans exit at a median ~40 months; defaults cluster 30–80 months."),
    sp(4),
    add_image(img("E1_eda_categorical.png"), max_h=2.6*inch),
    caption("Fig 2. Event breakdown by loan purpose and occupancy status. "
            "Cash-out Refi and Investment properties carry the highest default rates."),
    sp(4),
    add_image(img("E1_eda_vintage.png"), max_h=2.6*inch),
    caption("Fig 3. Vintage cohort analysis. 2006–2008 vintages show peak default rates "
            "of 8–12%; post-2012 cohorts exhibit near-zero default with high prepayment."),
    sp(4),
    add_image(img("E1_hazard_overall.png"), max_h=2.8*inch),
    caption("Fig 4. Cause-specific hazard intensity h(t). Prepayment hazard peaks early "
            "(months 10–30) then decays — the burnout effect. Default hazard rises slowly "
            "and plateaus, reflecting gradual financial deterioration."),
    sp(4),
    add_image(img("E1_hazard_stratified.png"), max_h=4.5*inch),
    caption("Fig 5. Stratified hazard intensity by FICO, LTV, origination rate, and vintage era. "
            "High-FICO loans peak earlier (faster refinancing); high-LTV loans show elevated "
            "default hazard throughout; GFC vintages (1999–2007) show a default hazard spike "
            "absent in post-GFC cohorts."),
    PageBreak(),
]

story += [h2("(b) Aalen-Johansen Cumulative Incidence Functions"), sp(2)]
story += [
    p("<b>Kaplan-Meier</b> (incorrect under competing risks):"),
    formula("S_KM(t) = prod_{t_i <= t} (1 - d_i^(k) / n_i)     F_KM^(k)(t) = 1 - S_KM(t)"),
    p("where d_i^(k) counts only cause-k events; competing exits treated as censored."),
    sp(4),
    p("<b>Aalen-Johansen</b> (correct):"),
    formula("F_k(t) = SUM_{t_j <= t}  S(t_j-)  *  d_{kj} / n_j"),
    Table([
        [Paragraph("<b>Symbol</b>", small_style), Paragraph("<b>Meaning</b>", small_style)],
        [Paragraph("F_k(t)", small_style), Paragraph("CIF for cause k at loan age t", small_style)],
        [Paragraph("d_{kj}", small_style), Paragraph("cause-k exits at time t_j", small_style)],
        [Paragraph("n_j", small_style),    Paragraph("loans at risk just before t_j", small_style)],
        [Paragraph("S(t_j-)", small_style),Paragraph("overall survival before t_j — uses ALL exits", small_style)],
    ], colWidths=[1.3*inch, 5.0*inch],
       style=TableStyle([
           ("BACKGROUND", (0,0), (-1,0), BLUE), ("TEXTCOLOR", (0,0), (-1,0), white),
           ("ROWBACKGROUNDS", (0,1), (-1,-1), [white, LGREY]),
           ("GRID", (0,0), (-1,-1), 0.5, HexColor("#bdc3c7")),
           ("FONTSIZE", (0,0), (-1,-1), 8.5),
           ("TOPPADDING", (0,0), (-1,-1), 3), ("BOTTOMPADDING", (0,0), (-1,-1), 3),
       ])),
    sp(4),
    p("<b>Partition identity</b> (impossible under KM):"),
    formula("SUM_k  F_k(t)  +  S(t)  =  1   for all t"),
    sp(4),
    add_image(img("E1_competing_risks_cif.png"), max_h=3.0*inch),
    caption("Fig 6. AJ cumulative incidence functions. "
            "10-yr prepayment CIF ≈ 83%;  default CIF ≈ 2.2%;  still active ≈ 14.5%."),
    PageBreak(),
]

story += [h2("(c) KM vs AJ — Bias from Ignoring Competing Events"), sp(2)]
story += [
    p("KM treats defaulted loans as <b>independent censored</b> observations. "
      "In practice they are the riskiest loans — least likely to prepay. "
      "Removing them inflates the estimated prepayment hazard at every subsequent time point."),
    sp(4),
    p("<b>KM increment</b> at event time t_j (defaults excluded from denominator):"),
    formula("ΔF_KM(t_j)  =  S_KM(t_j-)  *  d_{1j} / n_j^KM        [n_j^KM < n_j  →  over-estimates]"),
    sp(2),
    p("<b>AJ increment</b> at the same time (overall survival used):"),
    formula("ΔF_AJ(t_j)  =  S(t_j-)     *  d_{1j} / n_j            [S(t_j-) <= S_KM(t_j-)  always]"),
    sp(2),
    p("Because S(t_j-) ≤ S_KM(t_j-) at every step, every AJ increment is ≤ the KM increment, "
      "so the bias compounds over time:"),
    sp(2),
    b("5-year horizon:  KM over-estimates prepayment CIF by ~+0.45 pp"),
    b("10-year horizon: bias ≈ +1.31 pp"),
    b("20-year horizon: bias ≈ +2.23 pp"),
    sp(4),
    add_image(img("E1_km_vs_aj_bias.png"), max_h=2.8*inch),
    caption("Fig 7. KM vs AJ prepayment CIF — growing bias as competing exits accumulate."),
    PageBreak(),
]

story += [h2("(d) Stratified CIF by Loan Characteristics"), sp(2)]
story += [
    p("The prepayment/default balance varies substantially across loan segments. "
      "High-FICO borrowers show elevated prepayment CIF and near-zero default; "
      "high-LTV loans show flatter prepayment curves and elevated default risk; "
      "crisis vintages (2006–2008) display default rates an order of magnitude "
      "above modern cohorts."),
    sp(4),
    add_image(img("E1_stratified_cif.png"), max_h=3.4*inch),
    caption("Fig 6. Stratified AJ CIF by FICO bucket, LTV bucket, and vintage era."),
    PageBreak(),
]

story += [h2("(e) Cause-Specific Cox Regression  (12 features)"), sp(2)]
story += [
    p("Two separate Cox models are fitted — one for each cause — treating the other "
      "cause as censored at exit. Features: 8 numeric (FICO, LTV, DTI, UPB, "
      "mortgage rate, unemployment, HPI YoY, rate incentive) + 4 one-hot "
      "(loan purpose × 2, occupancy × 2). Both charts sorted ascending by log-HR."),
    sp(4),
    add_image(img("E1_cox_prepay.png"), max_h=3.4*inch),
    caption("Fig 8. Prepayment cause-specific Cox forest plot (sorted ascending). "
            "Rate incentive and LoanPurpose=Purchase drive the largest positive hazard; "
            "CreditScore has a positive effect (high-FICO refinances faster)."),
    sp(4),
    add_image(img("E1_cox_default.png"), max_h=3.4*inch),
    caption("Fig 9. Default cause-specific Cox forest plot (sorted ascending). "
            "LTV and DTI drive the largest positive default hazard; CreditScore strongly "
            "negative — opposite sign to prepayment, confirming the need for cause-specific models."),
    sp(4),
    p("<b>Key finding:</b> CreditScore and LTV have <b>opposite signs</b> between causes. "
      "A naive combined-Cox coefficient averages these out and misrepresents both risks."),
    PageBreak(),
]

story += [h2("(f) Fine-Gray Subdistribution Hazard"), sp(2)]
story += [
    p("The Fine-Gray model directly targets the CIF rather than the cause-specific "
      "hazard. Subjects who experienced a competing event (default) are kept in the "
      "risk set with IPCW weights, per Fine &amp; Gray (1999). "
      "The subdistribution hazard integrates to the CIF, so covariate effects are "
      "directly interpretable in terms of the observable event probability."),
    sp(4),
    Table([
        [Paragraph("<b>Model</b>", small_style),
         Paragraph("<b>Risk set at t</b>", small_style),
         Paragraph("<b>Coefficient interpretation</b>", small_style)],
        [Paragraph("Cause-specific Cox", small_style),
         Paragraph("All uncensored, non-defaulted", small_style),
         Paragraph("Effect on cause-specific hazard", small_style)],
        [Paragraph("<b>Fine-Gray</b>", small_style),
         Paragraph("All non-prepaid (incl. defaulted)", small_style),
         Paragraph("<b>Direct effect on CIF</b>", small_style)],
    ], colWidths=[1.8*inch, 2.4*inch, 2.4*inch],
       style=TableStyle([
           ("BACKGROUND", (0,0), (-1,0), BLUE),
           ("TEXTCOLOR",  (0,0), (-1,0), white),
           ("BACKGROUND", (0,2), (-1,2), LGREY),
           ("GRID", (0,0), (-1,-1), 0.5, HexColor("#bdc3c7")),
           ("FONTSIZE", (0,0), (-1,-1), 8.5),
           ("TOPPADDING",  (0,0), (-1,-1), 4),
           ("BOTTOMPADDING", (0,0), (-1,-1), 4),
       ])),
    sp(6),
    p("For prepayment, <b>rate_incentive</b> has a larger coefficient under Fine-Gray "
      "than cause-specific Cox: loans with high rate incentive are also less likely to "
      "default first, amplifying the competing-risk-adjusted CIF effect."),
    PageBreak(),
]

# ── E(ii): Time-Varying ───────────────────────────────────────────────────────
story += h1("E(ii) — Time-Dependent Covariates"); story += [sp(4)]
story += [
    p("Standard Cox models fix covariates at origination. The <b>Andersen-Gill "
      "counting-process</b> extension allows each loan to contribute one row per "
      "calendar month with that month's covariate values. The monthly panel dataset "
      "(93M rows, 2M loans) provides mortgage rate, unemployment, HPI YoY, and "
      "rate incentive updated each period."),
    sp(4),
    Table([
        [Paragraph("<b>Model</b>", small_style),
         Paragraph("<b>Covariate assumption</b>", small_style),
         Paragraph("<b>Data format</b>", small_style)],
        [Paragraph("Standard Cox — E(i)(e)", small_style),
         Paragraph("Fixed at origination vintage", small_style),
         Paragraph("One row per loan", small_style)],
        [Paragraph("<b>Andersen-Gill Cox — E(ii)</b>", small_style),
         Paragraph("Updated each month", small_style),
         Paragraph("One row per loan-month", small_style)],
    ], colWidths=[2.2*inch, 2.2*inch, 2.2*inch],
       style=TableStyle([
           ("BACKGROUND", (0,0), (-1,0), BLUE),
           ("TEXTCOLOR",  (0,0), (-1,0), white),
           ("BACKGROUND", (0,2), (-1,2), LGREY),
           ("GRID", (0,0), (-1,-1), 0.5, HexColor("#bdc3c7")),
           ("FONTSIZE", (0,0), (-1,-1), 8.5),
           ("TOPPADDING", (0,0), (-1,-1), 4),
           ("BOTTOMPADDING", (0,0), (-1,-1), 4),
       ])),
    sp(8),
    p("<b>Sample:</b> 10,000 loans stratified by vintage year, yielding ~500K "
      "loan-month rows. 11 features: 7 static (FICO, LTV, DTI, UPB, + 3 one-hot) + "
      "4 time-varying (mortgage rate, unemployment, HPI YoY, rate incentive)."),
    sp(6),
    p("<b>Key finding:</b> the current <b>rate_incentive</b> coefficient is larger "
      "in magnitude under the Andersen-Gill model than the origination-vintage "
      "snapshot. A loan originated when rates were high has a permanently large "
      "static rate_incentive; the TV model tracks the actual monthly refi "
      "option value and is more responsive to rate cycles. Unemployment and HPI "
      "also show stronger effects when measured contemporaneously."),
    sp(4),
]

for fname in ["Eii_tv_cox_forest.png", "Eii_static_vs_tv.png"]:
    p_img = img(fname)
    if p_img:
        story += [add_image(p_img, max_h=3.0*inch),
                  caption(f"Fig — {fname.replace('_',' ').replace('.png','')}"),
                  sp(4)]
    else:
        story += [p(f"<i>[{fname} — generated on notebook run]</i>"), sp(4)]

story.append(PageBreak())

# ── E(iv): Scenario Analysis ──────────────────────────────────────────────────
story += h1("E(iv) — Scenario Analysis: Interest Rate Shocks"); story += [sp(4)]
story += [
    p("Both Deep Cox and XGBoost are shocked with ±100 bp and ±200 bp changes to "
      "<i>mortgage_rate</i>. Rate incentive is updated consistently: "
      "rate_incentive = orig_rate − (mortgage_rate + shock)."),
    sp(6),
    add_image(img("E2_scenario_analysis.png"), max_h=3.2*inch),
    caption("Fig 8. Prepayment sensitivity to interest rate shocks. "
            "XGBoost (right) shows clear monotone convex response; "
            "Deep Cox (left) shows a smoother log-hazard ratio response."),
    sp(8),
    h2("Interpretation"),
    sp(2),
    b("XGBoost shows a clear monotone, convex response: rate cuts raise "
      "prepayment probability more than equivalent rate rises reduce it."),
    b("Deep Cox responds monotonically with smaller relative changes — the "
      "log-hazard is a smoother function of rate_incentive."),
    b("Asymmetry (rate cuts > rate rises) reflects the burnout effect: the most "
      "rate-sensitive borrowers have already prepaid, dampening upside response."),
    b("Both models agree on direction, confirming robustness of the rate-incentive "
      "channel across architectures."),
    PageBreak(),
]

# ── Summary table ─────────────────────────────────────────────────────────────
story += h1("Summary of Part E Extensions"); story += [sp(6)]
tbl_data = [
    [Paragraph("<b>Section</b>", small_style),
     Paragraph("<b>Method</b>", small_style),
     Paragraph("<b>Key Result</b>", small_style)],
    [Paragraph("E(i)(b) AJ CIF", small_style),
     Paragraph("Aalen-Johansen", small_style),
     Paragraph("10-yr prepay CIF ≈ 45%; default ≈ 2%", small_style)],
    [Paragraph("E(i)(c) KM Bias", small_style),
     Paragraph("KM vs AJ comparison", small_style),
     Paragraph("KM over-estimates by +0.5–2.2 pp at 5–20 yr", small_style)],
    [Paragraph("E(i)(d) Stratified CIF", small_style),
     Paragraph("AJ by FICO / LTV / era", small_style),
     Paragraph("Crisis vintages show 8–12× higher default CIF", small_style)],
    [Paragraph("E(i)(e) Cause-Specific Cox", small_style),
     Paragraph("Two CoxPH models, 12 features", small_style),
     Paragraph("CreditScore: opposite signs across causes", small_style)],
    [Paragraph("E(i)(f) Fine-Gray", small_style),
     Paragraph("IPCW-weighted Cox", small_style),
     Paragraph("rate_incentive HR larger vs cause-specific", small_style)],
    [Paragraph("E(ii) TV Cox", small_style),
     Paragraph("Andersen-Gill, 11 features", small_style),
     Paragraph("Current rate_incentive > origination snapshot", small_style)],
    [Paragraph("E(iv) Scenarios", small_style),
     Paragraph("Deep Cox + XGBoost shock", small_style),
     Paragraph("Convex, asymmetric prepayment response", small_style)],
]
tbl = Table(tbl_data, colWidths=[1.6*inch, 1.9*inch, 3.1*inch])
tbl.setStyle(TableStyle([
    ("BACKGROUND", (0,0), (-1,0), BLUE),
    ("TEXTCOLOR",  (0,0), (-1,0), white),
    ("BACKGROUND", (0,1), (-1,-1), white),
    ("ROWBACKGROUNDS", (0,1), (-1,-1), [white, LGREY]),
    ("GRID", (0,0), (-1,-1), 0.5, HexColor("#bdc3c7")),
    ("FONTSIZE", (0,0), (-1,-1), 8.5),
    ("TOPPADDING", (0,0), (-1,-1), 5),
    ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
]))
story.append(tbl)

# ── Build PDF ─────────────────────────────────────────────────────────────────
pdf_path = str(OUT_DIR / "PartE_Report.pdf")
doc = SimpleDocTemplate(pdf_path, pagesize=letter,
                        leftMargin=0.9*inch, rightMargin=0.9*inch,
                        topMargin=0.8*inch, bottomMargin=0.8*inch)
doc.build(story)
print(f"PDF saved → {pdf_path}")

# ─────────────────────────────────────────────────────────────────────────────
# 2. POWERPOINT PRESENTATION  (python-pptx)
# ─────────────────────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Madrid theme
DARK  = RGBColor(0x1F, 0x39, 0x63)   # deep navy
LIGHT = RGBColor(0xC0, 0x00, 0x00)   # crimson red
GOLD_C= RGBColor(0xFF, 0xC0, 0x00)   # gold accent stripe
TEAL_C= RGBColor(0x00, 0x56, 0x7A)   # teal for secondary
LGRY  = RGBColor(0xF2, 0xF2, 0xF2)
MGRY  = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_C = RGBColor(0xC0, 0x00, 0x00)

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h, size=20, bold=False, color=None,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb

def add_img(slide, path, l, t, w, h=None):
    if not path:
        return
    try:
        with PILImage.open(path) as im:
            iw, ih = im.size
        if h is None:
            h = w * ih / iw
        max_h = 5.8
        if h > max_h:
            w = w * max_h / h
            h = max_h
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    except Exception as e:
        print(f"  img error {path}: {e}")

def header_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.05, DARK)          # navy header
    rect(slide, 0, 1.05, 13.33, 0.07, GOLD_C)     # gold accent stripe
    txbox(slide, title, 0.3, 0.06, 12.5, 0.55, size=26, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.3, 0.60, 12.5, 0.4, size=13, color=LGRY)

def formula_box(slide, lines, l, t, w, h, size=13):
    """Monospaced formula text box with light-blue tinted background."""
    bg = RGBColor(0xEE, 0xF2, 0xF7)
    rect(slide, l, t, w, h, fill_rgb=bg)
    tb = slide.shapes.add_textbox(Inches(l+0.1), Inches(t+0.08),
                                   Inches(w-0.2), Inches(h-0.1))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = para.add_run(); run.text = line
        run.font.size = Pt(size)
        run.font.name = "Courier New"
        run.font.color.rgb = DARK

def bullet_box(slide, items, l, t, w, h, size=14):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            para = tf.paragraphs[0]; first = False
        else:
            para = tf.add_paragraph()
        para.text = f"• {item}"
        para.font.size = Pt(size)
        para.space_after = Pt(4)

# ── Slide 1: Title  (Madrid theme) ───────────────────────────────────────────
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, DARK)           # navy background
rect(sl, 0, 2.5, 13.33, 0.1, GOLD_C)       # gold horizontal divider
rect(sl, 0, 2.6, 13.33, 2.5, LIGHT)        # crimson band
rect(sl, 0, 5.1, 13.33, 0.1, GOLD_C)       # gold horizontal divider
txbox(sl, "MTH9877 — Interest Rates & Credit", 1, 0.9, 11.33, 0.7,
      size=18, color=LGRY, align=PP_ALIGN.CENTER)
txbox(sl, "Assignment 3 — Part E Extensions", 1, 1.6, 11.33, 0.85,
      size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Mortgage Prepayment & Default  |  Survival Analysis Extensions",
      1, 2.72, 11.33, 0.65, size=20, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Rose Lin  ·  Baruch MFE  ·  May 2026",
      1, 6.55, 11.33, 0.5, size=13, color=LGRY, align=PP_ALIGN.CENTER)

# ── Slide 2: Outline ──────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Part E — Overview", "Four extension tracks")
rect(sl, 0, 1.1, 13.33, 6.4, LGRY)
items = [
    "E(i)   Competing Risks Framework — Aalen-Johansen CIF, KM bias, stratified CIF, cause-specific Cox, Fine-Gray",
    "E(ii)  Time-Dependent Covariates — Andersen-Gill counting-process Cox with monthly macro",
    "E(iii) Neural Survival Models — DeepHit framework (architecture overview)",
    "E(iv)  Scenario Analysis — Prepayment sensitivity to ±200 bp rate shocks",
]
bullet_box(sl, items, 0.5, 1.3, 12.3, 5.5, size=17)

# ── Slide 3: EDA — duration & covariates ─────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(a) — EDA: Duration Distributions & Covariate Profiles",
           "Prepaid: median ~40 mo  |  Defaulted: median ~55 mo  |  full 34M-loan dataset")
add_img(sl, img("E1_eda_duration.png"), 0.15, 1.15, 13.0)

# ── Slide 4: EDA — vintage & categorical ─────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(a) — EDA: Vintage Cohort & Categorical Breakdown")
add_img(sl, img("E1_eda_vintage.png"), 0.15, 1.15, 8.0)
add_img(sl, img("E1_eda_categorical.png"), 8.3, 1.15, 4.8)

# ── Slide 5: EDA — overall hazard intensity ──────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(a) — Hazard Intensity: Prepayment Peaks Early, Default Rises Slowly",
           "Cause-specific Nelson-Aalen smooth · burnout visible in prepayment hazard · 100K sample")
add_img(sl, img("E1_hazard_overall.png"), 0.5, 1.2, 12.3)

# ── Slide 6: EDA — stratified hazard ─────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(a) — Stratified Hazard Intensity by FICO · LTV · Rate · Vintage",
           "High-FICO refinances fastest · High-LTV defaults persist longer · GFC vintages spike")
add_img(sl, img("E1_hazard_stratified.png"), 0.5, 1.1, 12.3)

# ── Slide 6: AJ CIF — formula slide ──────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(b) — Aalen-Johansen CIF: Formula & Contrast",
           "AJ correctly weights each increment by the overall survival probability")
rect(sl, 0, 1.12, 13.33, 6.38, LGRY)

# KM box
txbox(sl, "Kaplan-Meier  (INCORRECT — treats competing exits as censored)", 0.3, 1.25, 12.7, 0.4,
      size=13, bold=True, color=LIGHT)
formula_box(sl,
    ["S_KM(t) = PROD_{t_i <= t} (1 - d_i^(k) / n_i)",
     "F_KM^(k)(t) = 1 - S_KM(t)    ← overcounts cause-k hazard"],
    0.3, 1.65, 12.7, 0.9, size=13)

# AJ box
txbox(sl, "Aalen-Johansen  (CORRECT)", 0.3, 2.65, 12.7, 0.4,
      size=13, bold=True, color=DARK)
formula_box(sl,
    ["F_k(t) = SUM_{t_j <= t}  S(t_j-)  *  d_{kj} / n_j",
     "",
     "  F_k(t) : CIF for cause k at loan age t",
     "  d_{kj} : cause-k exits at time t_j",
     "  n_j    : loans at risk just before t_j",
     "  S(t_j-): overall survival before t_j — uses ALL exits"],
    0.3, 3.05, 12.7, 1.55, size=13)

# Partition identity
txbox(sl, "Partition identity  (impossible under KM):", 0.3, 4.68, 12.7, 0.35,
      size=12, bold=True, color=DARK)
formula_box(sl, ["SUM_k  F_k(t)  +  S(t)  =  1    for all t"],
            0.3, 5.05, 12.7, 0.5, size=13)

# Chart
add_img(sl, img("E1_competing_risks_cif.png"), 0.3, 5.6, 12.7)

# ── Slide 6b: AJ CIF chart full ───────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(b) — AJ CIF: 10-yr Prepay 83%  ·  Default 2.2%  ·  Active 14.5%",
           "S_KM over-estimates each cause; AJ partition sums exactly to 1")
add_img(sl, img("E1_competing_risks_cif.png"), 0.5, 1.2, 12.3)

# ── Slide 7: KM vs AJ — formula slide ────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(c) — KM vs AJ: Why Defaults Cannot Be Treated as Censorings",
           "Every AJ increment ≤ KM increment  →  F_AJ(t) ≤ F_KM(t)  →  growing bias")
rect(sl, 0, 1.12, 13.33, 6.38, LGRY)

txbox(sl, "KM increment at t_j  (defaults excluded from denominator n_j^KM < n_j):",
      0.3, 1.25, 12.7, 0.4, size=12, bold=True, color=LIGHT)
formula_box(sl,
    ["ΔF_KM(t_j)  =  S_KM(t_j-)  *  d_{1j} / n_j^KM    ← n_j^KM underestimates risk set"],
    0.3, 1.65, 12.7, 0.6, size=13)

txbox(sl, "AJ increment at t_j  (overall survival, full risk set):",
      0.3, 2.35, 12.7, 0.4, size=12, bold=True, color=DARK)
formula_box(sl,
    ["ΔF_AJ(t_j)   =  S(t_j-)     *  d_{1j} / n_j      ← S(t_j-) ≤ S_KM(t_j-) always"],
    0.3, 2.75, 12.7, 0.6, size=13)

txbox(sl, "Consequence — bias compounds over time:", 0.3, 3.45, 12.7, 0.4,
      size=12, bold=True, color=DARK)
bullet_box(sl, ["5 yr:  KM over-estimates prepayment CIF by ~+0.45 pp",
                "10 yr: bias ≈ +1.31 pp",
                "20 yr: bias ≈ +2.23 pp  (gap widens as defaulted loans accumulate)"],
           0.3, 3.85, 12.7, 1.1, size=13)
add_img(sl, img("E1_km_vs_aj_bias.png"), 0.3, 5.0, 12.7)

# ── Slide 8: Stratified CIF ──────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(d) — Stratified CIF",
           "High-FICO → elevated prepayment, near-zero default;  high-LTV → inverted")
add_img(sl, img("E1_stratified_cif.png"), 0.3, 1.15, 12.7)

# ── Slide 9a: Cause-Specific Cox — Prepayment ────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(e) — Cause-Specific Cox: Prepayment  (sorted ascending)",
           "Rate incentive & Purchase purpose → highest refi hazard  ·  12 features  ·  95% CI")
add_img(sl, img("E1_cox_prepay.png"), 0.3, 1.15, 12.7)

# ── Slide 9b: Cause-Specific Cox — Default ───────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(e) — Cause-Specific Cox: Default  (sorted ascending)",
           "LTV & DTI → highest default hazard  ·  CreditScore opposite sign to prepayment")
add_img(sl, img("E1_cox_default.png"), 0.3, 1.15, 12.7)

# ── Slide 10: Fine-Gray ───────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(i)(f) — Fine-Gray Subdistribution Hazard",
           "IPCW-weighted Cox approximation; coefficients directly target the CIF")
rect(sl, 0, 1.1, 13.33, 6.4, LGRY)
items = [
    "Subjects who experienced a competing event (default) are KEPT in the risk set with IPCW weights",
    "The subdistribution hazard integrates to the CIF — covariate effects are CIF effects",
    "rate_incentive HR is LARGER under Fine-Gray than cause-specific Cox:",
    "   high-incentive loans are also less likely to default first → amplifies CIF effect",
    "Low-FICO Fine-Gray HR for prepayment is ATTENUATED: elevated default risk reduces",
    "   measured prepayment propensity once competing risks are accounted for",
    "Use Fine-Gray for CIF prediction / loss forecasting",
    "Use cause-specific Cox for economic mechanism estimation",
]
bullet_box(sl, items, 0.5, 1.3, 12.3, 5.5, size=14)

# ── Slide 11: E(ii) TV Cox ────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(ii) — Time-Dependent Covariates",
           "Andersen-Gill counting-process Cox — monthly rate_incentive, unemployment, HPI")
rect(sl, 0, 1.1, 13.33, 6.4, LGRY)
items = [
    "Panel dataset: 93M rows × 2M loans — mortgage_rate, unemployment, HPI YoY, rate_incentive updated monthly",
    "10,000 loans stratified by vintage → ~500K counting-process rows (tstart, tstop, event)",
    "11 features: 7 static (FICO, LTV, DTI, UPB + 3 one-hot) + 4 time-varying",
    "Key finding: current rate_incentive coefficient LARGER than origination-vintage snapshot",
    "   Static Cox: rate_incentive locked at origination — high-rate era loans always look attractive",
    "   TV Cox: tracks actual monthly refi option value — more responsive to rate cycles",
    "Unemployment: contemporaneous job loss predicts default better than vintage average",
    "Limitation: origination LTV used (not ELTV) — adding current LTV would improve default hazard",
]
bullet_box(sl, items, 0.5, 1.3, 12.3, 5.5, size=13)

# ── Slide 12: E(ii) plots (if available) ─────────────────────────────────────
tv_f = img("Eii_tv_cox_forest.png")
tv_c = img("Eii_static_vs_tv.png")
if tv_f or tv_c:
    sl = add_slide()
    header_bar(sl, "E(ii) — TV Cox Results")
    if tv_f and tv_c:
        add_img(sl, tv_f, 0.15, 1.15, 7.8)
        add_img(sl, tv_c, 8.1, 1.15, 5.0)
    elif tv_f:
        add_img(sl, tv_f, 0.5, 1.15, 12.3)
    else:
        add_img(sl, tv_c, 0.5, 1.15, 12.3)

# ── Slide 13: E(iv) Scenario Analysis ────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(iv) — Scenario Analysis: Interest Rate Shocks",
           "±100 bp and ±200 bp shocks to mortgage_rate  (rate_incentive updated consistently)")
add_img(sl, img("E2_scenario_analysis.png"), 0.5, 1.15, 12.3)

# ── Slide 14: E(iv) Interpretation ───────────────────────────────────────────
sl = add_slide()
header_bar(sl, "E(iv) — Scenario Analysis: Interpretation")
rect(sl, 0, 1.1, 13.33, 6.4, LGRY)
items = [
    "XGBoost: clear monotone CONVEX response — rate cuts raise prepayment probability more than rate rises suppress it",
    "Deep Cox: monotone but smoother — log-hazard is a less non-linear function of rate_incentive",
    "Asymmetry (−200 bp effect > +200 bp effect) reflects the BURNOUT effect:",
    "   most rate-sensitive borrowers have already prepaid → muted upside response to cuts",
    "Both models agree on direction → robustness of rate-incentive channel across architectures",
    "Implication for MBS pricing: prepayment option has POSITIVE CONVEXITY —",
    "   rate decline creates more prepayment than rate rise suppresses it",
]
bullet_box(sl, items, 0.5, 1.3, 12.3, 5.5, size=15)

# ── Slide 15: Summary ─────────────────────────────────────────────────────────
sl = add_slide()
header_bar(sl, "Summary")
rect(sl, 0, 1.1, 13.33, 6.4, LGRY)
rows = [
    ("E(i)(b)", "Aalen-Johansen CIF", "10-yr prepay ≈ 45%;  default ≈ 2%"),
    ("E(i)(c)", "KM vs AJ bias", "KM over-estimates by +0.5–2.2 pp at 5–20 yr"),
    ("E(i)(d)", "Stratified CIF", "Crisis vintages 8–12× higher default CIF"),
    ("E(i)(e)", "Cause-specific Cox", "CreditScore / LTV opposite signs across causes"),
    ("E(i)(f)", "Fine-Gray", "rate_incentive HR larger vs cause-specific Cox"),
    ("E(ii)",   "Andersen-Gill TV Cox", "Current rate_incentive > origination snapshot"),
    ("E(iv)",   "Scenario analysis", "Convex, asymmetric prepayment response to rate shocks"),
]
y = 1.35
for sec, method, result in rows:
    rect(sl, 0.3, y, 1.2, 0.52, DARK)
    txbox(sl, sec, 0.32, y+0.05, 1.15, 0.42, size=12, bold=True, color=WHITE)
    txbox(sl, method, 1.6, y+0.05, 2.8, 0.42, size=12, bold=True, color=DARK)
    txbox(sl, result, 4.5, y+0.05, 8.5, 0.42, size=12, color=DARK)
    y += 0.60

pptx_path = str(OUT_DIR / "PartE_Presentation.pptx")
prs.save(pptx_path)
print(f"PPTX saved → {pptx_path}")
